"""
pptx_exporter.py — Exportador de relatórios para PPTX via python-pptx.

Dependência opcional: pip install python-pptx
Se python-pptx não estiver instalado, o exporter emite um aviso e retorna None.

Estratégia de Slides:
- Slide 1: Título e data
- Slide 2: Resumo Executivo (primeira seção H2 encontrada)
- Slide 3 em diante: Um slide por item H3 (projetos/ferramentas)
- Último slide: Recomendação Final
"""
import logging
import re
from pathlib import Path
from typing import Optional, List, Tuple

logger = logging.getLogger("pptx_exporter")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False
    logger.warning(
        "PPTXExporter: python-pptx não está instalado. "
        "Execute `pip install python-pptx` para habilitar exportação PPTX."
    )

_BRAND_BG = (15, 23, 42)       # Azul escuro
_BRAND_ACCENT = (99, 102, 241)  # Índigo
_BRAND_TEXT = (248, 250, 252)   # Quase branco


class PPTXExporter:
    """
    Converte um relatório Markdown (string) em uma apresentação PPTX estruturada usando python-pptx.
    Opera em modo degradado (retorna None com warning) se python-pptx não estiver disponível.
    """

    def __init__(self):
        self.available = _PPTX_AVAILABLE

    def export(self, markdown_content: str, filepath: str) -> Optional[str]:
        """
        Gera o PPTX a partir do conteúdo Markdown e o salva em `filepath`.

        Returns:
            Caminho do arquivo gerado, ou None se python-pptx não disponível.
        """
        if not self.available:
            logger.warning("PPTXExporter: exportação ignorada — python-pptx não disponível.")
            return None

        pptx_path = str(filepath).replace(".md", ".pptx")
        Path(pptx_path).parent.mkdir(parents=True, exist_ok=True)

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        sections = self._parse_sections(markdown_content)
        blank_layout = prs.slide_layouts[6]  # Blank

        for i, (heading, level, body_lines) in enumerate(sections):
            slide = prs.slides.add_slide(blank_layout)
            self._set_bg(slide, _BRAND_BG)

            if level == 0:
                # Slide de título
                self._add_text_box(
                    slide, heading, Inches(0.5), Inches(2.5), Inches(12), Inches(1.5),
                    font_size=Pt(40), bold=True, color=_BRAND_TEXT
                )
                if body_lines:
                    self._add_text_box(
                        slide, "\n".join(body_lines[:2]), Inches(0.5), Inches(4.2),
                        Inches(12), Inches(1.5),
                        font_size=Pt(18), bold=False, color=(160, 170, 200)
                    )
            elif level == 2:
                # Slide de seção principal (H2)
                self._add_accent_bar(slide)
                self._add_text_box(
                    slide, heading, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                    font_size=Pt(28), bold=True, color=_BRAND_TEXT
                )
                body_text = self._clean(
                    "\n".join(l for l in body_lines[:8] if l.strip())
                )
                self._add_text_box(
                    slide, body_text, Inches(0.5), Inches(1.3), Inches(12), Inches(5.5),
                    font_size=Pt(14), bold=False, color=_BRAND_TEXT
                )
            else:
                # Slide de item (H3+)
                self._add_accent_bar(slide)
                self._add_text_box(
                    slide, self._clean(heading), Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                    font_size=Pt(22), bold=True, color=_BRAND_TEXT
                )
                body_text = self._clean(
                    "\n".join(l.strip() for l in body_lines[:10] if l.strip())
                )
                self._add_text_box(
                    slide, body_text, Inches(0.5), Inches(1.3), Inches(12), Inches(5.5),
                    font_size=Pt(13), bold=False, color=_BRAND_TEXT
                )

        try:
            prs.save(pptx_path)
            logger.info(f"PPTXExporter: PPTX gerado em {pptx_path} ({len(prs.slides)} slides)")
            return pptx_path
        except Exception as e:
            logger.error(f"PPTXExporter: falha ao gerar PPTX: {e}")
            return None

    # ── Helpers ───────────────────────────────────────────────────────────────

    def _parse_sections(self, content: str) -> List[Tuple[str, int, List[str]]]:
        """
        Divide o Markdown em seções [(heading, level, body_lines)].
        level 0 = título (H1), 2 = H2, 3 = H3+
        """
        sections: List[Tuple[str, int, List[str]]] = []
        current_heading = ""
        current_level = 0
        current_body: List[str] = []

        for line in content.splitlines():
            stripped = line.strip()
            if stripped.startswith("# ") and not stripped.startswith("## "):
                if current_heading or current_body:
                    sections.append((current_heading, current_level, current_body))
                current_heading = stripped[2:]
                current_level = 0
                current_body = []
            elif stripped.startswith("## ") and not stripped.startswith("### "):
                if current_heading or current_body:
                    sections.append((current_heading, current_level, current_body))
                current_heading = stripped[3:]
                current_level = 2
                current_body = []
            elif stripped.startswith("### "):
                if current_heading or current_body:
                    sections.append((current_heading, current_level, current_body))
                current_heading = stripped[4:]
                current_level = 3
                current_body = []
            else:
                current_body.append(line)

        if current_heading or current_body:
            sections.append((current_heading, current_level, current_body))

        return sections[:30]  # máximo de 30 slides

    def _set_bg(self, slide, color: Tuple[int, int, int]) -> None:
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*color)

    def _add_accent_bar(self, slide) -> None:
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0), Inches(0), Inches(13.33), Inches(0.07)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*_BRAND_ACCENT)
        shape.line.fill.background()

    def _add_text_box(self, slide, text: str, left, top, width, height,
                      font_size=None, bold=False,
                      color: Tuple[int, int, int] = (255, 255, 255)) -> None:
        from pptx.dml.color import RGBColor
        from pptx.util import Pt as _Pt
        if font_size is None:
            font_size = _Pt(14)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        run = p.runs[0] if p.runs else p.add_run()
        run.text = text
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color)

    def _clean(self, text: str) -> str:
        """
        Remove Markdown básico para texto limpo nos slides.
        """
        text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
        text = re.sub(r"\*(.+?)\*", r"\1", text)
        text = re.sub(r"`(.+?)`", r"\1", text)
        text = re.sub(r"\[([^\]]+)\]\(([^\)]+)\)", r"\1", text)
        text = re.sub(r"^#+\s+", "", text, flags=re.MULTILINE)
        text = re.sub(r"^>\s+", "", text, flags=re.MULTILINE)
        text = re.sub(r"^-\s+", "• ", text, flags=re.MULTILINE)
        text = re.sub(r"^\d+\.\s+", "• ", text, flags=re.MULTILINE)
        text = re.sub(r"---+", "", text)
        return text.strip()
